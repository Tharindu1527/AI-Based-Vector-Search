import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import os

class DocumentLoader:
    """Extract text from various document formats"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF using PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX using python-docx"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from DOCX: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX using python-pptx"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise Exception(f"Error extracting text from PPTX: {str(e)}")
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                raise Exception(f"Error reading TXT file with multiple encodings: {str(e)}")
        except Exception as e:
            raise Exception(f"Error extracting text from TXT: {str(e)}")
    
    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """Extract text based on file extension"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return cls.extract_text_from_pdf(file_path)
        elif file_extension == '.docx':
            return cls.extract_text_from_docx(file_path)
        elif file_extension == '.pptx':
            return cls.extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            return cls.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    @staticmethod
    def get_supported_extensions():
        """Return list of supported file extensions"""
        return ['.pdf', '.docx', '.pptx', '.txt']
    
    @classmethod
    def is_supported_file(cls, filename: str) -> bool:
        """Check if file extension is supported"""
        file_extension = os.path.splitext(filename)[1].lower()
        return file_extension in cls.get_supported_extensions()
    
    @staticmethod
    def get_file_info(file_path: str) -> dict:
        """Get basic information about the file"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_stats = os.stat(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        
        return {
            'filename': filename,
            'extension': file_extension,
            'size_bytes': file_stats.st_size,
            'size_mb': round(file_stats.st_size / (1024 * 1024), 2),
            'is_supported': file_extension in DocumentLoader.get_supported_extensions()
        }